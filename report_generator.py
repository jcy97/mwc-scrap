from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.colors import Color, blue, black, darkblue
from reportlab.pdfbase import pdfutils
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime


class PDFGenerator:
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.setup_styles()

    def setup_styles(self):
        self.styles.add(
            ParagraphStyle(
                name="CustomTitle",
                parent=self.styles["Title"],
                fontSize=24,
                spaceAfter=30,
                textColor=darkblue,
                alignment=1,
            )
        )

        self.styles.add(
            ParagraphStyle(
                name="CategoryTitle",
                parent=self.styles["Heading1"],
                fontSize=16,
                spaceAfter=12,
                textColor=blue,
                leftIndent=20,
            )
        )

        self.styles.add(
            ParagraphStyle(
                name="RuleItem",
                parent=self.styles["Normal"],
                fontSize=10,
                spaceAfter=8,
                leftIndent=40,
                rightIndent=20,
                textColor=black,
            )
        )

    def generate_pdf(self, data):
        filename = f"exhibition_rules_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        doc = SimpleDocTemplate(filename, pagesize=A4)
        story = []

        title = Paragraph("전시회 규칙 및 규정 가이드", self.styles["CustomTitle"])
        story.append(title)
        story.append(Spacer(1, 20))

        subtitle = Paragraph(
            f"생성일: {datetime.now().strftime('%Y년 %m월 %d일')}",
            self.styles["Normal"],
        )
        story.append(subtitle)
        story.append(Spacer(1, 30))

        for rule_data in data:
            rule_title = Paragraph(f"📋 {rule_data['title']}", self.styles["Heading1"])
            story.append(rule_title)
            story.append(Spacer(1, 12))

            url_para = Paragraph(
                f"출처: <link href='{rule_data['url']}'>{rule_data['url']}</link>",
                self.styles["Normal"],
            )
            story.append(url_para)
            story.append(Spacer(1, 20))

            for category in rule_data["categories"]:
                cat_title = Paragraph(
                    f"🔸 {category['name']}", self.styles["CategoryTitle"]
                )
                story.append(cat_title)
                story.append(Spacer(1, 8))

                for idx, item in enumerate(category["items"], 1):
                    if len(item) > 20:
                        item_para = Paragraph(f"{idx}. {item}", self.styles["RuleItem"])
                        story.append(item_para)

                story.append(Spacer(1, 15))

            story.append(PageBreak())

        doc.build(story)
        return filename


class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slides()

    def setup_slides(self):
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "전시회 규칙 및 규정 가이드"
        subtitle.text = f"Exhibition Rules & Regulations\n생성일: {datetime.now().strftime('%Y년 %m월 %d일')}"

    def generate_ppt(self, data):
        for rule_data in data:
            title_slide = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(title_slide)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = rule_data["title"]

            tf = content.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"출처: {rule_data['url']}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(128, 128, 128)

            for category in rule_data["categories"]:
                content_slide = self.prs.slide_layouts[1]
                slide = self.prs.slides.add_slide(content_slide)
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = f"🔸 {category['name']}"

                tf = content.text_frame
                tf.clear()

                for idx, item in enumerate(category["items"][:8], 1):
                    if len(item) > 15:
                        p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
                        p.text = (
                            f"{idx}. {item[:150]}{'...' if len(item) > 150 else ''}"
                        )
                        p.font.size = Pt(14)
                        p.space_after = Pt(8)

        summary_slide = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(summary_slide)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "📊 요약"

        tf = content.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"• 총 {len(data)}개의 규칙 문서"
        p.font.size = Pt(16)

        total_categories = sum(len(rule["categories"]) for rule in data)
        p = tf.add_paragraph()
        p.text = f"• 총 {total_categories}개의 카테고리"
        p.font.size = Pt(16)

        p = tf.add_paragraph()
        p.text = f"• 생성일: {datetime.now().strftime('%Y년 %m월 %d일 %H:%M')}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(128, 128, 128)

        filename = f"exhibition_rules_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        self.prs.save(filename)
        return filename
